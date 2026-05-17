import pytest
from pathlib import Path
from unittest.mock import patch, MagicMock
from reader import read_file, ReadResult

def test_read_eml(tmp_dir):
    eml_content = (
        b"From: alice@example.com\r\n"
        b"To: bob@example.com\r\n"
        b"Subject: Case Update\r\n"
        b"Date: Mon, 1 Jan 2024 10:00:00 +0000\r\n"
        b"Content-Type: text/plain\r\n\r\n"
        b"This is the email body."
    )
    eml_path = tmp_dir / "test.eml"
    eml_path.write_bytes(eml_content)

    result = read_file(eml_path)
    assert isinstance(result, ReadResult)
    assert "From: alice@example.com" in result.text
    assert "This is the email body." in result.text
    assert result.extension == ".eml"
    assert result.ocr is False

def test_read_docx(tmp_dir):
    from docx import Document
    doc = Document()
    doc.add_paragraph("Introduction paragraph.")
    doc.add_paragraph("Second paragraph about the case.")
    docx_path = tmp_dir / "test.docx"
    doc.save(str(docx_path))

    result = read_file(docx_path)
    assert "Introduction paragraph." in result.text
    assert "Second paragraph about the case." in result.text
    assert result.extension == ".docx"
    assert result.ocr is False

def test_read_pptx(tmp_dir):
    from pptx import Presentation
    prs = Presentation()
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Legal Overview"
    slide.placeholders[1].text = "Key points about the statute."
    pptx_path = tmp_dir / "test.pptx"
    prs.save(str(pptx_path))

    result = read_file(pptx_path)
    assert "Legal Overview" in result.text
    assert result.extension == ".pptx"
    assert result.ocr is False

def test_read_msg(tmp_dir):
    mock_msg = MagicMock()
    mock_msg.sender = "alice@example.com"
    mock_msg.to = "bob@example.com"
    mock_msg.subject = "Update"
    mock_msg.date = "2024-01-01"
    mock_msg.body = "MSG body text."

    msg_path = tmp_dir / "test.msg"
    msg_path.write_bytes(b"fake msg content")

    with patch("reader.extract_msg_lib.Message", return_value=mock_msg):
        result = read_file(msg_path)

    assert "alice@example.com" in result.text
    assert "MSG body text." in result.text
    assert result.extension == ".msg"

def test_unsupported_extension_raises(tmp_dir):
    bad_path = tmp_dir / "file.xyz"
    bad_path.write_text("content")
    with pytest.raises(ValueError, match="Unsupported format"):
        read_file(bad_path)

def test_read_result_has_required_fields(tmp_dir):
    eml_path = tmp_dir / "test.eml"
    eml_path.write_bytes(
        b"From: x@x.com\r\nSubject: Test\r\nContent-Type: text/plain\r\n\r\nBody."
    )
    result = read_file(eml_path)
    assert hasattr(result, 'text')
    assert hasattr(result, 'filename')
    assert hasattr(result, 'extension')
    assert hasattr(result, 'size_bytes')
    assert hasattr(result, 'ocr')
    assert hasattr(result, 'ocr_confidence')
