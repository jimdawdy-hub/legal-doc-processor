import email as email_lib
import subprocess
import tempfile
from dataclasses import dataclass
from pathlib import Path
from typing import Optional

from pdfminer.high_level import extract_text as pdfminer_extract

try:
    import fitz  # pymupdf
    FITZ_AVAILABLE = True
except ImportError:
    FITZ_AVAILABLE = False

try:
    import pytesseract
    from PIL import Image
    TESSERACT_AVAILABLE = True
except ImportError:
    TESSERACT_AVAILABLE = False

import extract_msg as extract_msg_lib


@dataclass
class ReadResult:
    text: str
    filename: str
    extension: str
    size_bytes: int
    page_count: Optional[int] = None
    ocr: bool = False
    ocr_confidence: Optional[float] = None


SUPPORTED_EXTENSIONS = {'.pdf', '.docx', '.pptx', '.eml', '.msg'}


def read_file(path: Path) -> Optional['ReadResult']:
    """
    Returns ReadResult, or None if the file is a scanned PDF with OCR
    confidence below threshold (caller should move to ocr_queue).
    """
    ext = path.suffix.lower()
    if ext not in SUPPORTED_EXTENSIONS:
        raise ValueError(f"Unsupported format: {ext}")
    size = path.stat().st_size

    if ext == '.pdf':
        return _read_pdf(path, size)
    elif ext == '.docx':
        return _read_docx(path, size)
    elif ext == '.pptx':
        return _read_pptx(path, size)
    elif ext == '.eml':
        return _read_eml(path, size)
    elif ext == '.msg':
        return _read_msg(path, size)


def _read_docx(path: Path, size: int) -> ReadResult:
    from docx import Document
    doc = Document(str(path))
    parts = []
    for para in doc.paragraphs:
        if para.text.strip():
            parts.append(para.text)
    for table in doc.tables:
        for row in table.rows:
            row_text = ' | '.join(c.text.strip() for c in row.cells if c.text.strip())
            if row_text:
                parts.append(row_text)
    return ReadResult(
        text='\n'.join(parts), filename=path.name,
        extension='.docx', size_bytes=size,
    )


def _read_pptx(path: Path, size: int) -> ReadResult:
    from pptx import Presentation
    prs = Presentation(str(path))
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        parts.append(f"[Slide {i}]")
        for shape in slide.shapes:
            if hasattr(shape, 'text') and shape.text.strip():
                parts.append(shape.text.strip())
        if (slide.has_notes_slide and
                slide.notes_slide.notes_text_frame.text.strip()):
            parts.append(f"[Notes] {slide.notes_slide.notes_text_frame.text.strip()}")
    return ReadResult(
        text='\n'.join(parts), filename=path.name,
        extension='.pptx', size_bytes=size,
    )


def _read_eml(path: Path, size: int) -> ReadResult:
    with open(path, 'rb') as f:
        msg = email_lib.message_from_bytes(f.read())

    headers = {k: msg.get(k, '') for k in ('From', 'To', 'Subject', 'Date', 'CC', 'BCC')}
    header_text = '\n'.join(f"{k}: {v}" for k, v in headers.items() if v)

    body = ''
    if msg.is_multipart():
        for part in msg.walk():
            if part.get_content_type() == 'text/plain':
                payload = part.get_payload(decode=True)
                if payload:
                    body = payload.decode('utf-8', errors='replace')
                    break
    else:
        payload = msg.get_payload(decode=True)
        if payload:
            body = payload.decode('utf-8', errors='replace')

    return ReadResult(
        text=f"{header_text}\n\n{body}".strip(),
        filename=path.name, extension='.eml', size_bytes=size,
    )


def _read_msg(path: Path, size: int) -> ReadResult:
    msg = extract_msg_lib.Message(str(path))
    headers = (
        f"From: {msg.sender or ''}\n"
        f"To: {msg.to or ''}\n"
        f"Subject: {msg.subject or ''}\n"
        f"Date: {msg.date or ''}"
    )
    return ReadResult(
        text=f"{headers}\n\n{msg.body or ''}".strip(),
        filename=path.name, extension='.msg', size_bytes=size,
    )


def _read_pdf(path: Path, size: int) -> Optional[ReadResult]:
    text = pdfminer_extract(str(path)) or ''
    if len(text.strip()) >= 100:
        return ReadResult(
            text=text, filename=path.name,
            extension='.pdf', size_bytes=size, ocr=False,
        )
    return _ocr_pdf(path, size)


def _ocr_pdf(path: Path, size: int) -> Optional[ReadResult]:
    """Returns None if OCR confidence is below 70% threshold."""
    confidence = _sample_ocr_confidence(path)
    if confidence is not None and confidence < 70.0:
        return None  # caller moves to ocr_queue

    with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as tmp:
        tmp_path = Path(tmp.name)
    try:
        subprocess.run(
            ['ocrmypdf', '--output-type', 'pdf', '--redo-ocr',
             str(path), str(tmp_path)],
            check=True, capture_output=True,
        )
        text = pdfminer_extract(str(tmp_path)) or ''
    finally:
        tmp_path.unlink(missing_ok=True)

    return ReadResult(
        text=text, filename=path.name, extension='.pdf',
        size_bytes=size, ocr=True, ocr_confidence=confidence,
    )


def _sample_ocr_confidence(path: Path) -> Optional[float]:
    """Sample first 3 pages with pytesseract; return mean word confidence or None."""
    if not (FITZ_AVAILABLE and TESSERACT_AVAILABLE):
        return None

    doc = fitz.open(str(path))
    confidences = []
    for page_num in range(min(3, len(doc))):
        pix = doc[page_num].get_pixmap(dpi=200)
        img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
        data = pytesseract.image_to_data(img, output_type=pytesseract.Output.DICT)
        page_conf = [
            c for c in data['conf']
            if isinstance(c, (int, float)) and c > 0
        ]
        confidences.extend(page_conf)

    return sum(confidences) / len(confidences) if confidences else None
